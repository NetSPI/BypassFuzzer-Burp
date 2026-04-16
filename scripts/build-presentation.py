#!/usr/bin/env python3
"""Build the BypassFuzzer team-pitch presentation.

Dark theme, 16:9, monospace for code snippets. Attribution on every
slide that reuses third-party research (Tiurin, ZeroNights 2018).

Run from repo root:
    python3 scripts/build-presentation.py
Writes to docs/BypassFuzzer-Presentation.pptx.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BG            = RGBColor(0x10, 0x12, 0x18)  # near-black
FG            = RGBColor(0xE5, 0xE7, 0xEB)  # light grey
ACCENT        = RGBColor(0xF4, 0xA8, 0x1E)  # amber for emphasis
DIM           = RGBColor(0x9C, 0xA3, 0xAF)  # muted
CODE_BG       = RGBColor(0x1B, 0x1F, 0x27)  # slightly lighter than bg
CODE_FG       = RGBColor(0xCB, 0xE7, 0xFF)  # soft blue
ATTR          = RGBColor(0x6B, 0x72, 0x80)  # dim grey for attribution
GOOD          = RGBColor(0x4A, 0xDE, 0x80)  # green
BAD           = RGBColor(0xF4, 0x71, 0x71)  # red

TITLE_FONT    = "Inter"
BODY_FONT     = "Inter"
MONO_FONT     = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def paint_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=FG, font=BODY_FONT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    tb.text_frame.margin_left = Inches(0.05)
    tb.text_frame.margin_right = Inches(0.05)
    tb.text_frame.margin_top = 0
    tb.text_frame.margin_bottom = 0
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_code_block(slide, left, top, width, height, lines, *, size=16,
                   color=CODE_FG):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CODE_BG
    card.line.color.rgb = RGBColor(0x30, 0x35, 0x3E)
    card.line.width = Pt(0.75)
    card.shadow.inherit = False

    tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2),
                                   width - Inches(0.5), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, (text, line_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name = MONO_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = line_color if line_color else color


def add_attribution(slide, text):
    add_text(
        slide,
        Inches(0.5), Inches(7.0),
        Inches(12.33), Inches(0.35),
        text,
        size=10, color=ATTR, font=BODY_FONT,
    )


def add_footer_accent(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(6.85),
        Inches(0.4), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# Slide 1 — title
s1 = prs.slides.add_slide(blank)
paint_background(s1, BG)

add_text(s1, Inches(0.75), Inches(2.6), Inches(11.8), Inches(1.2),
         "BypassFuzzer", size=64, bold=True, color=ACCENT, font=TITLE_FONT)
add_text(s1, Inches(0.75), Inches(3.55), Inches(11.8), Inches(0.9),
         "Automating AuthZ Bypass Testing", size=28, color=FG, font=TITLE_FONT)
add_text(s1, Inches(0.75), Inches(4.5), Inches(11.8), Inches(0.5),
         "Burp Suite extension · Bypass · URL Validation · IDOR tabs",
         size=16, color=DIM, font=BODY_FONT)

# Presenter + date line (date intentionally open — edit in-deck)
add_text(s1, Inches(0.75), Inches(6.3), Inches(11.8), Inches(0.4),
         "Jonathan Conesa", size=18, bold=True, color=FG, font=BODY_FONT)
add_text(s1, Inches(0.75), Inches(6.75), Inches(11.8), Inches(0.35),
         "Security Consultant  ·  NetSpi", size=12, color=ATTR, font=BODY_FONT)

# accent bar
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.75), Inches(2.3),
    Inches(0.5), Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background(); bar.shadow.inherit = False


# Slide 2 — references / inspiration
s_ref = prs.slides.add_slide(blank)
paint_background(s_ref, BG)

add_text(s_ref, Inches(0.75), Inches(0.55), Inches(11.8), Inches(0.7),
         "Built on a decade of public research.", size=34, bold=True, color=FG, font=TITLE_FONT)
add_text(s_ref, Inches(0.75), Inches(1.25), Inches(11.8), Inches(0.5),
         "Path and URL parser bugs have been a drumbeat since 2018. Three talks/resources we lean on heavily:",
         size=16, color=DIM, font=BODY_FONT)

# Three reference cards stacked vertically
def reference_card(slide, top_in, venue, title, author_line, url):
    top = Inches(top_in)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), top, Inches(11.8), Inches(1.35))
    card.fill.solid(); card.fill.fore_color.rgb = CODE_BG
    card.line.color.rgb = RGBColor(0x30, 0x35, 0x3E); card.line.width = Pt(0.75)
    card.shadow.inherit = False

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.75), top, Inches(0.08), Inches(1.35))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background(); accent.shadow.inherit = False

    add_text(slide, Inches(1.05), top + Inches(0.1), Inches(11.5), Inches(0.3),
             venue, size=11, bold=True, color=ACCENT, font=BODY_FONT)
    add_text(slide, Inches(1.05), top + Inches(0.35), Inches(11.5), Inches(0.45),
             title, size=18, bold=True, color=FG, font=TITLE_FONT)
    add_text(slide, Inches(1.05), top + Inches(0.75), Inches(11.5), Inches(0.3),
             author_line, size=12, color=DIM, font=BODY_FONT)
    add_text(slide, Inches(1.05), top + Inches(1.02), Inches(11.5), Inches(0.3),
             url, size=10, color=CODE_FG, font=MONO_FONT)

reference_card(
    s_ref, 2.05,
    "BLACK HAT USA 2018 · DEF CON 26",
    "Breaking Parser Logic! Take Your Path Normalization Off and Pop 0days Out!",
    "Orange Tsai  ·  @orange_8361  ·  DEVCORE",
    "i.blackhat.com/us-18/Wed-August-8/us-18-Orange-Tsai-Breaking-Parser-Logic-...pdf",
)

reference_card(
    s_ref, 3.55,
    "ZERONIGHTS 2018",
    "Reverse Proxies & Inconsistency",
    "Aleksei \"GreenDog\" Tiurin  ·  @antyurin  ·  Acunetix",
    "2018.zeronights.ru/wp-content/uploads/materials/20-Reverse-proxies-Inconsistency.pdf",
)

reference_card(
    s_ref, 5.05,
    "PORTSWIGGER RESEARCH  (ongoing)",
    "URL Validation Bypass Cheat Sheet",
    "PortSwigger  ·  open community PRs",
    "portswigger.net/web-security/ssrf/url-validation-bypass-cheat-sheet",
)

# Framing
add_text(s_ref, Inches(0.75), Inches(6.7), Inches(11.8), Inches(0.4),
         "Plus PoCs and CVE write-ups from Bentkowski, Kettle, Wallarm, and many more.",
         size=13, color=DIM, font=BODY_FONT)

add_footer_accent(s_ref)


# Slide 3 — what we're bypassing (the defender's rule)
s2 = prs.slides.add_slide(blank)
paint_background(s2, BG)

add_text(s2, Inches(0.75), Inches(0.55), Inches(11.8), Inches(0.7),
         "What we're trying to bypass.", size=36, bold=True, color=FG, font=TITLE_FONT)
add_text(s2, Inches(0.75), Inches(1.25), Inches(11.8), Inches(0.5),
         "Every stack ships a flavor of \"block requests to /admin.\" Teams trust these rules.",
         size=18, color=DIM, font=BODY_FONT)

# Three protection patterns, side by side
col_w = Inches(4.0)
col_gap = Inches(0.25)
left0 = Inches(0.45)
top_labels = Inches(2.15)
top_codes = Inches(2.55)

# Column 1 — Nginx
add_text(s2, left0, top_labels, col_w, Inches(0.4),
         "Nginx", size=15, bold=True, color=ACCENT, font=BODY_FONT)
add_code_block(s2, left0, top_codes, col_w, Inches(1.8), [
    ("location /admin {",               CODE_FG),
    ("    deny all;",                    CODE_FG),
    ("    return 403;",                  CODE_FG),
    ("}",                                 CODE_FG),
], size=14)

# Column 2 — Spring Security
left1 = left0 + col_w + col_gap
add_text(s2, left1, top_labels, col_w, Inches(0.4),
         "Spring Security", size=15, bold=True, color=ACCENT, font=BODY_FONT)
add_code_block(s2, left1, top_codes, col_w, Inches(1.8), [
    ("http.authorizeRequests()",           CODE_FG),
    ("  .antMatchers(\"/admin/**\")",      CODE_FG),
    ("    .hasRole(\"ADMIN\")",            CODE_FG),
    ("  .anyRequest().permitAll();",       CODE_FG),
], size=14)

# Column 3 — HAProxy
left2 = left1 + col_w + col_gap
add_text(s2, left2, top_labels, col_w, Inches(0.4),
         "HAProxy", size=15, bold=True, color=ACCENT, font=BODY_FONT)
add_code_block(s2, left2, top_codes, col_w, Inches(1.8), [
    ("acl restricted path_beg /admin",   CODE_FG),
    ("http-request deny if restricted",   CODE_FG),
    ("",                                   CODE_FG),
    ("# blocks anything starting /admin", DIM),
], size=14)

# Takeaway
box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.75), Inches(4.85), Inches(11.8), Inches(1.6))
box.fill.solid(); box.fill.fore_color.rgb = CODE_BG
box.line.color.rgb = ACCENT; box.line.width = Pt(1.5)
box.shadow.inherit = False

add_text(s2, Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.5),
         "All three rules compare against the literal string \"/admin\".",
         size=18, bold=True, color=FG, font=BODY_FONT)
add_text(s2, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.5),
         "None of them account for every way a client can reach the /admin handler.",
         size=18, color=FG, font=BODY_FONT)
add_text(s2, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.4),
         "The URL parser that matches, and the URL parser that routes, are not the same parser.",
         size=14, color=DIM, font=BODY_FONT)

add_footer_accent(s2)


# Slide 3 — Not a one-off
s3 = prs.slides.add_slide(blank)
paint_background(s3, BG)

add_text(s3, Inches(0.75), Inches(0.55), Inches(11.8), Inches(0.7),
         "Not a one-off. Every stack has these.", size=36, bold=True, color=FG, font=TITLE_FONT)
add_text(s3, Inches(0.75), Inches(1.25), Inches(11.8), Inches(0.5),
         "Three more from the same talk. Each is a real bypass, in a real stack.",
         size=18, color=DIM, font=BODY_FONT)

# Example 1
add_text(s3, Inches(0.75), Inches(2.0), Inches(11.8), Inches(0.4),
         "Nginx + Weblogic", size=16, bold=True, color=ACCENT, font=BODY_FONT)
add_code_block(s3, Inches(0.75), Inches(2.45), Inches(11.8), Inches(1.0), [
    ("GET /#/../Login.jsp HTTP/1.1", CODE_FG),
    ("     Nginx sees:     /        (# is a fragment)", DIM),
    ("     Weblogic sees:  /Login.jsp   ← bypasses the block", GOOD),
], size=14)

# Example 2
add_text(s3, Inches(0.75), Inches(3.65), Inches(11.8), Inches(0.4),
         "Nginx + Tomcat", size=16, bold=True, color=ACCENT, font=BODY_FONT)
add_code_block(s3, Inches(0.75), Inches(4.1), Inches(11.8), Inches(1.0), [
    ("GET /iframe_safe/..;/admin HTTP/1.1", CODE_FG),
    ("     Nginx sees:     /iframe_safe/... (inside the allow rule)", DIM),
    ("     Tomcat sees:    /admin   ← matrix param ;/ collapses", GOOD),
], size=14)

# Example 3
add_text(s3, Inches(0.75), Inches(5.3), Inches(11.8), Inches(0.4),
         "HAProxy path_beg /admin", size=16, bold=True, color=ACCENT, font=BODY_FONT)
add_code_block(s3, Inches(0.75), Inches(5.75), Inches(11.8), Inches(1.0), [
    ("GET /%61dmin HTTP/1.1", CODE_FG),
    ("     HAProxy sees:   /%61dmin (literal, no match)", DIM),
    ("     Origin sees:    /admin   ← URL-decoded before dispatch", GOOD),
], size=14)

add_footer_accent(s3)
add_attribution(s3, "All three examples recreated from: Aleksei Tiurin — \"Reverse Proxies & Inconsistency,\" ZeroNights 2018")


# Slide 4 — Why this doesn't scale manually
s4 = prs.slides.add_slide(blank)
paint_background(s4, BG)

add_text(s4, Inches(0.75), Inches(0.55), Inches(11.8), Inches(0.7),
         "Doing this by hand doesn't scale.", size=36, bold=True, color=FG, font=TITLE_FONT)
add_text(s4, Inches(0.75), Inches(1.25), Inches(11.8), Inches(0.5),
         "For a single /admin returning 403, a thorough test has to try:",
         size=18, color=DIM, font=BODY_FONT)

# Bullet-style list as pseudo-code
bullets = [
    ("~15 encoding tricks", "single, double, triple, overlong UTF-8, fullwidth Unicode"),
    ("~8 path-shape primitives", "../  ..;/  ..%2f  /./  sacrificial-prefix  matrix-param  ..\\  //"),
    ("~5 whitespace/control insertions", "%09  %0a  %0d  %20  %00"),
    ("cross-combinations of the above", "../..%2f   ..;/%2e%2e%2f   %2e%09%2e   /x/../  /%u002e  ..."),
    ("target-preservation constraint", "every variant must still resolve to /admin after server-side normalization"),
]
top = Inches(2.1)
for label, detail in bullets:
    add_text(s4, Inches(0.75), top, Inches(4.2), Inches(0.45),
             "• " + label, size=17, bold=True, color=FG, font=BODY_FONT)
    add_text(s4, Inches(4.9), top, Inches(7.65), Inches(0.45),
             detail, size=14, color=DIM, font=MONO_FONT)
    top += Inches(0.55)

# Takeaway box
box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.75), Inches(5.3), Inches(11.8), Inches(1.4))
box.fill.solid(); box.fill.fore_color.rgb = CODE_BG
box.line.color.rgb = ACCENT; box.line.width = Pt(1.5)
box.shadow.inherit = False

add_text(s4, Inches(1.0), Inches(5.45), Inches(11.3), Inches(0.5),
         "Manual fuzz of one endpoint can take forever, maybe 50 variants tried before you give up.",
         size=18, bold=True, color=FG, font=BODY_FONT)
add_text(s4, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.5),
         "BypassFuzzer Path attack against /admin: ~14,000 variants, ~45 seconds.",
         size=18, bold=True, color=ACCENT, font=BODY_FONT)
add_text(s4, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.35),
         "Every one of them target-preserving — a 200 means you bypassed AuthZ on THAT endpoint.",
         size=13, color=DIM, font=BODY_FONT)

add_footer_accent(s4)


# ---------------------------------------------------------------------------

repo_root = Path(__file__).resolve().parent.parent
out_path = repo_root / "docs" / "BypassFuzzer-Presentation.pptx"
out_path.parent.mkdir(exist_ok=True)
prs.save(out_path)
print(f"wrote {out_path} ({out_path.stat().st_size} bytes, {len(prs.slides)} slides)")
